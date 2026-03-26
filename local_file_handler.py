"""
Local File Handler
Reads documents from local folder instead of Google Drive
Supports: PDF, Word, Excel, PowerPoint, CSV, TXT, and more
"""

import os
from typing import List, Dict, Tuple, Optional
import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation
import pandas as pd
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor, as_completed

# Safety limits — override via environment variables if needed
MAX_FILE_SIZE_MB = float(os.getenv('KB_MAX_FILE_SIZE_MB', '50'))
MAX_PDF_PAGES    = int(os.getenv('KB_MAX_PDF_PAGES', '300'))
MAX_EXCEL_ROWS   = int(os.getenv('KB_MAX_EXCEL_ROWS', '5000'))


class LocalFileHandler:
    def __init__(self, base_folder: str = "kb_raw"):
        # Resolve to absolute path once so all comparisons are consistent
        self.base_folder = os.path.abspath(base_folder)

        if not os.path.exists(self.base_folder):
            print(f"⚠️  Warning: Folder '{base_folder}' does not exist")
            print(f"   Creating folder: {base_folder}")
            os.makedirs(self.base_folder)
        else:
            print(f"✓ Found knowledge base folder: {base_folder}")

    def fetch_all_documents(self, folder_path: str = None) -> List[Dict]:
        """
        Recursively fetch all documents from local folder.
        Uses parallel I/O so multiple files are extracted simultaneously.

        Args:
            folder_path: Optional subfolder path (relative to base_folder).
                         Paths that escape base_folder are silently rejected.
        """
        if folder_path is None:
            target = self.base_folder
        else:
            # Security: realpath resolves ".." and symlinks; ensure result stays
            # inside base_folder before we touch it.
            candidate = os.path.realpath(os.path.join(self.base_folder, folder_path))
            base_with_sep = self.base_folder.rstrip(os.sep) + os.sep
            if not (candidate == self.base_folder or candidate.startswith(base_with_sep)):
                print(f"⚠️  Rejected folder_path outside knowledge base: {folder_path}")
                target = self.base_folder
            else:
                target = candidate

        # Phase 1: collect all file paths (fast, sequential directory walk)
        file_entries: List[Tuple[str, str, str]] = []
        self._collect_paths(target, file_entries)

        if not file_entries:
            print(f"\n✓ Found 0 documents")
            return []

        # Phase 2: extract content in parallel (bottleneck is disk I/O + PDF/Office parsing)
        workers = min(8, len(file_entries))
        all_documents: List[Dict] = []
        with ThreadPoolExecutor(max_workers=workers) as executor:
            futures = {
                executor.submit(self._process_one_file, ip, iname, rp): rp
                for ip, iname, rp in file_entries
            }
            for future in as_completed(futures):
                try:
                    result = future.result()
                    if result:
                        all_documents.append(result)
                except Exception as e:
                    print(f"  ✗ Unexpected error: {str(e)}")

        # Sort for deterministic ordering (parallel completion is non-deterministic)
        all_documents.sort(key=lambda d: d['path'])

        print(f"\n✓ Found {len(all_documents)} documents")
        return all_documents

    # ── Internal helpers ───────────────────────────────────

    def _collect_paths(self, folder_path: str, entries: List[Tuple], relative_path: str = ""):
        """Walk folder tree and collect (item_path, item_name, rel_path) — no extraction."""
        try:
            for item in sorted(os.listdir(folder_path)):
                item_path = os.path.join(folder_path, item)
                current_path = os.path.join(relative_path, item) if relative_path else item

                if item.startswith('.') or item == '__MACOSX':
                    continue

                # Never follow symlinks — prevents path traversal via symlinks in kb_raw
                if os.path.islink(item_path):
                    continue

                if os.path.isdir(item_path):
                    print(f"📁 Scanning folder: {current_path}")
                    self._collect_paths(item_path, entries, current_path)
                else:
                    entries.append((item_path, item, current_path))
        except Exception as e:
            print(f"Error scanning folder {folder_path}: {str(e)}")

    def _process_one_file(self, item_path: str, item_name: str, rel_path: str) -> Optional[Dict]:
        """Extract content from a single file with a size guard. Returns dict or None."""
        try:
            size_mb = os.path.getsize(item_path) / (1024 * 1024)
            if size_mb > MAX_FILE_SIZE_MB:
                print(f"  ⚠  Skipping {item_name}: {size_mb:.1f} MB > {MAX_FILE_SIZE_MB} MB limit")
                return None
        except OSError:
            return None

        print(f"📄 Processing: {rel_path}")
        doc_content = self._extract_content(item_path, item_name)

        if doc_content:
            return {
                'id': item_path,
                'name': item_name,
                'path': rel_path,
                'content': doc_content,
                'mimeType': self._get_mime_type(item_name),
                'modifiedTime': os.path.getmtime(item_path),
            }
        return None

    # ── MIME helper ────────────────────────────────────────

    def _get_mime_type(self, filename: str) -> str:
        ext = os.path.splitext(filename)[1].lower()
        mime_types = {
            '.pdf':  'application/pdf',
            '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            '.doc':  'application/msword',
            '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            '.xls':  'application/vnd.ms-excel',
            '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            '.ppt':  'application/vnd.ms-powerpoint',
            '.txt':  'text/plain',
            '.csv':  'text/csv',
            '.md':   'text/markdown',
            '.json': 'application/json',
        }
        return mime_types.get(ext, 'application/octet-stream')

    # ── Content extraction ─────────────────────────────────

    def _extract_content(self, file_path: str, filename: str) -> str:
        ext = os.path.splitext(filename)[1].lower()
        try:
            if ext == '.pdf':
                return self._extract_pdf(file_path)
            elif ext in ['.docx', '.doc']:
                return self._extract_docx(file_path)
            elif ext in ['.xlsx', '.xls']:
                return self._extract_excel(file_path)
            elif ext in ['.pptx', '.ppt']:
                return self._extract_pptx(file_path)
            elif ext == '.csv':
                return self._extract_csv(file_path)
            elif ext in ['.txt', '.md', '.json', '.log']:
                return self._extract_text(file_path)
            else:
                print(f"  ⚠  Unsupported file type: {ext}")
                return ""
        except Exception as e:
            print(f"  ✗ Error extracting content: {str(e)}")
            return ""

    def _extract_pdf(self, file_path: str) -> str:
        try:
            text = ""
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                total_pages = len(pdf_reader.pages)
                pages_to_read = pdf_reader.pages[:MAX_PDF_PAGES]
                if total_pages > MAX_PDF_PAGES:
                    print(f"  ℹ  Large PDF: reading first {MAX_PDF_PAGES} of {total_pages} pages")
                for page in pages_to_read:
                    page_text = page.extract_text() or ""
                    text += page_text + "\n"
            return text
        except Exception as e:
            print(f"  ✗ PDF extraction error: {str(e)}")
            return ""

    def _extract_docx(self, file_path: str) -> str:
        try:
            doc = DocxDocument(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + " "
                    text += "\n"
            return text
        except Exception as e:
            print(f"  ✗ Word extraction error: {str(e)}")
            return ""

    def _extract_excel(self, file_path: str) -> str:
        try:
            excel_file = pd.ExcelFile(file_path)
            text = ""
            for sheet_name in excel_file.sheet_names:
                # nrows limit prevents loading huge sheets into memory
                df = pd.read_excel(file_path, sheet_name=sheet_name, nrows=MAX_EXCEL_ROWS)
                text += f"\n--- Sheet: {sheet_name} ---\n"
                text += df.to_string(index=False) + "\n"
            return text
        except Exception as e:
            print(f"  ✗ Excel extraction error: {str(e)}")
            return ""

    def _extract_pptx(self, file_path: str) -> str:
        try:
            prs = Presentation(file_path)
            text = ""
            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"\n--- Slide {slide_num} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            print(f"  ✗ PowerPoint extraction error: {str(e)}")
            return ""

    def _extract_csv(self, file_path: str) -> str:
        try:
            df = pd.read_csv(file_path)
            return df.to_string(index=False)
        except Exception as e:
            print(f"  ✗ CSV extraction error: {str(e)}")
            return ""

    def _extract_text(self, file_path: str) -> str:
        # Try UTF-8 first (the common case); fall back to latin-1, which never
        # raises UnicodeDecodeError (every byte is a valid latin-1 character).
        for encoding in ('utf-8', 'latin-1'):
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
        print(f"  ✗ Could not decode text file")
        return ""

    # ── Public utilities ───────────────────────────────────

    def check_connection(self) -> bool:
        return os.path.exists(self.base_folder) and os.path.isdir(self.base_folder)

    def get_folder_stats(self) -> Dict:
        stats = {'total_files': 0, 'total_folders': 0, 'file_types': {}, 'total_size_mb': 0}
        for root, dirs, files in os.walk(self.base_folder):
            # Don't follow symlinks in stats walk either
            dirs[:] = [d for d in dirs if not os.path.islink(os.path.join(root, d))]
            stats['total_folders'] += len(dirs)
            for file in files:
                if file.startswith('.'):
                    continue
                stats['total_files'] += 1
                ext = os.path.splitext(file)[1].lower()
                stats['file_types'][ext] = stats['file_types'].get(ext, 0) + 1
                file_path = os.path.join(root, file)
                stats['total_size_mb'] += os.path.getsize(file_path) / (1024 * 1024)
        return stats


if __name__ == "__main__":
    print("Testing Local File Handler")
    print("=" * 50)

    handler = LocalFileHandler("kb_raw")

    if handler.check_connection():
        print("\n✓ Folder is accessible")

        stats = handler.get_folder_stats()
        print(f"\nFolder Statistics:")
        print(f"  Total files: {stats['total_files']}")
        print(f"  Total folders: {stats['total_folders']}")
        print(f"  Total size: {stats['total_size_mb']:.2f} MB")
        print(f"  File types: {dict(stats['file_types'])}")

        print("\n" + "=" * 50)
        print("Fetching documents...")
        docs = handler.fetch_all_documents()

        if docs:
            print(f"\n✓ Successfully processed {len(docs)} documents")
            print(f"\nSample document:")
            print(f"  Name: {docs[0]['name']}")
            print(f"  Path: {docs[0]['path']}")
            print(f"  Content length: {len(docs[0]['content'])} characters")
            print(f"  Preview: {docs[0]['content'][:200]}...")
        else:
            print("\n⚠  No documents found. Add some files to kb_raw folder!")
    else:
        print("\n✗ Folder not accessible")
